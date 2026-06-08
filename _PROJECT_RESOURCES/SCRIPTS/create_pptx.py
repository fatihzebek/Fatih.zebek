from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

# Paths
folder_path = r'c:\Users\FatihZebek\Desktop\Dh_Servis\src\presentation\TUREK_Conference'
output_file = os.path.join(folder_path, 'TUREK_2025_Demirer_Holding.pptx')

prs = Presentation()

def add_slide(title_text, bullet_points=None, image_name=None):
    # Use blank layout to have full control
    slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(slide_layout)
    
    # Add Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(32)
    p.alignment = PP_ALIGN.CENTER

    # Add Image if provided
    if image_name:
        img_path = os.path.join(folder_path, image_name)
        if os.path.exists(img_path):
            # Place image on the left or bottom depending on content
            if bullet_points:
                slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), height=Inches(4.5))
            else:
                slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

    # Add Bullet Points
    if bullet_points:
        left = Inches(5.5) if image_name else Inches(1)
        top = Inches(1.5)
        width = Inches(4) if image_name else Inches(8)
        height = Inches(5)
        
        txBox2 = slide.shapes.add_textbox(left, top, width, height)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        
        for point in bullet_points:
            p = tf2.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(18)

# Slide 1: Cover
add_slide("Rüzgarın İkinci Baharı: Demirer Holding", 
          ["Türkiye'nin İlk Rüzgar Mirasından Geleceğin Enerji Vizyonuna", "TÜREK 2025 Konferansı"],
          "turek_title_slide_1778088108720.png")

# Slide 2: Öncü Miras
add_slide("Öncü Miras ve Değişen Gündem", 
          ["Türkiye'nin ilk RES işletmecisi olarak 25 yıllık tecrübe.", 
           "İlk nesil santrallerin ömür sonu yönetimi.",
           "Neden Repowering? (Verimlilik, Modernizasyon, Sürdürülebilirlik)"])

# Slide 3: Filo Analizi
add_slide("Mevcut Filo Analizi ve Stratejik Güç", 
          ["Toplam Türbin Sayısı: 284 Adet",
           "Toplam Kurulu Güç: ~368 MW",
           "Alize Enerji: 96 Türbin / 169 MW",
           "Anemon Enerji: 49 Türbin / 55.7 MW",
           "Doğal Enerji: 48 Türbin / 57.2 MW",
           "Mare Manastır: 55 Türbin / 56.2 MW",
           "Dares: 36 Türbin / 29.4 MW"],
          "fleet_analysis_infographic_1778088430993.png")

# Slide 4: Repowering
add_slide("Repowering - Kapasite ve Verim Artışı", 
          ["Küçük türbinlerden devasa rotorlara geçiş.",
           "Daha az kule, daha yüksek enerji yoğunluğu.",
           "Mevzuat ve lisanslama süreçleri."],
          "repowering_concept_1778088127723.png")

# Slide 5: Kış Operasyonları
add_slide("Kış Operasyonları ve YAT Talimatları", 
          ["Kritik Sorun: Yük Atma (YAT) talimatlarının fiziksel etkisi.",
           "Duran kanatlarda kontrolsüz buz birikimi.",
           "Üretim kayıpları ve restart zorlukları."],
          "winter_icing_turbine_1778088150477.png")

# Slide 6: Dijital O&M
add_slide("Dijital O&M ve Siber Güvenlik", 
          ["AI tabanlı öngörücü bakım.",
           "SCADA sistemlerinin siber zırhı.",
           "Akıllı rüzgar sahaları yönetimi."],
          "cybersecurity_wind_farm_1778088172356.png")

# Slide 7: Sonuç
add_slide("Sonuç ve Vizyon", 
          ["Demirer Holding olarak sektöre yeniden yön verme.",
           "Teşekkür ederiz.",
           "Soru & Cevap"])

prs.save(output_file)
print(f"PowerPoint file created at: {output_file}")

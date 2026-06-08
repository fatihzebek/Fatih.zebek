from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

# Paths
folder_path = r'c:\Users\FatihZebek\Desktop\Dh_Servis\src\presentation\TUREK_Conference'
output_file = os.path.join(folder_path, 'TUREK_2025_Demirer_Enercon_DeepAnalysis.pptx')

prs = Presentation()

def add_standard_slide(title_text, bullet_points=None, image_name=None):
    # Use Title and Content layout (standard index 1)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = title_text
    
    # Content Placeholder
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    
    if bullet_points:
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0

    # Add image as a separate shape if exists
    if image_name:
        img_path = os.path.join(folder_path, image_name)
        if os.path.exists(img_path):
            # Scale image to fit on the side or bottom
            slide.shapes.add_picture(img_path, Inches(6), Inches(2), width=Inches(3.5))

# 1. Kapak
slide_layout = prs.slide_layouts[0] # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "DEMİRER HOLDİNG & ENERCON"
subtitle.text = "Türkiye'nin Rüzgar Mirasından Geleceğin Teknolojisine\nTÜREK 2025 Konferansı"

# 2. Tarihsel Kilometre Taşları
add_standard_slide(
    "1996'dan Bugüne: Bir Başarı Hikayesi", 
    [
        "1996: İlk rüzgar ölçüm ağı (108 kule).",
        "1998: Germiyan RES - Türkiye'nin İLK rüzgar santrali.",
        "2000: Bozcaada RES - Enercon E-40 teknolojisi.",
        "25+ Yıllık Kesintisiz İşletme Tecrübesi."
    ],
    "enercon_legacy_genesis_1778089178519.png"
)

# 3. Enercon Teknolojisi
add_standard_slide(
    "Enercon: Dişlisiz (Gearless) Teknolojinin Gücü", 
    [
        "Direct Drive Teknolojisi: Minimum arıza, maksimum verimlilik.",
        "Aşınmasız Çalışma: Şanzıman kaynaklı risklerin sıfırlanması.",
        "Düşük Bakım Maliyeti: Uzun ömürlülüğün sırrı.",
        "Şebeke Uyumu: En zorlu koşullarda kararlı üretim."
    ],
    "enercon_direct_drive_tech_1778089196015.png"
)

# 4. Filo Analizi
add_standard_slide(
    "Filo Portföyü (355 Ünitelik Enercon Ağı)", 
    [
        "Toplam 355 Adet Enercon Türbini.",
        "103 Adet Enercon E-44 / 85 Adet Enercon E-48.",
        "Yeni Nesil: E-126, E-115, E-92 Serileri.",
        "Toplam Kurulu Güç: ~1 GW (Saha Verileri)."
    ],
    "fleet_analysis_infographic_1778088430993.png"
)

# 5. Kış Operasyonları
add_standard_slide(
    "Kış Operasyonları ve YAT Yönetimi", 
    [
        "YAT Talimatlarının donma riskine etkisi.",
        "Enercon Buz Algılama ve Isıtma Çözümleri.",
        "Sensör verileriyle güvenli restart stratejileri.",
        "Operasyonel sürekliliğin korunması."
    ],
    "enercon_winter_grid_safety_1778089233070.png"
)

# 6. Repowering
add_standard_slide(
    "Repowering: Rüzgarın İkinci Baharı", 
    [
        "E-40 ve E-44 sahalarında dönüşüm vizyonu.",
        "Aynı sahada 3 katına varan kapasite artışı.",
        "Yeni nesil E-138 / E-160 serisi geçişi.",
        "Sürdürülebilirlik ve döngüsel ekonomi."
    ],
    "enercon_evolution_repowering_1778089214937.png"
)

# 7. Sonuç
add_standard_slide(
    "Teşekkürler", 
    [
        "Demirer Holding & Enercon: 25 Yıllık Güçlü Ortaklık.",
        "Türkiye'nin enerji bağımsızlığına katkı.",
        "Soru & Cevap"
    ]
)

prs.save(output_file)
print(f"Standard PowerPoint file created at: {output_file}")

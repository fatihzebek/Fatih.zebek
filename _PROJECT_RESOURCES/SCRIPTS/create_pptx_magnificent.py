from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Paths
folder_path = r'c:\Users\FatihZebek\Desktop\Dh_Servis\src\presentation\TUREK_Conference'
output_file = os.path.join(folder_path, 'TUREK_2025_Demirer_Enercon_Final.pptx')

prs = Presentation()

def apply_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_magnificent_slide(title_text, subtitle_text=None, bullet_points=None, image_name=None, is_title_slide=False):
    slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark Theme
    apply_background(slide, RGBColor(10, 25, 45)) # Dark Navy

    # Title Styling
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(36) if not is_title_slide else Pt(44)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT if not is_title_slide else PP_ALIGN.CENTER

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0, 180, 216) # Teal
        p2.alignment = PP_ALIGN.LEFT if not is_title_slide else PP_ALIGN.CENTER

    # Image Placement
    if image_name:
        img_path = os.path.join(folder_path, image_name)
        if os.path.exists(img_path):
            if is_title_slide:
                slide.shapes.add_picture(img_path, Inches(1), Inches(2.2), width=Inches(8))
            else:
                slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.8), height=Inches(4.5))

    # Bullet Points with Premium Styling
    if bullet_points:
        left = Inches(5.5) if image_name else Inches(1)
        top = Inches(1.8)
        width = Inches(4) if image_name else Inches(8)
        
        txBox = slide.shapes.add_textbox(left, top, width, Inches(5))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = "• " + point
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(220, 220, 220)
            p.space_after = Pt(12)

# 1. Kapak
add_magnificent_slide(
    "DEMİRER HOLDİNG & ENERCON", 
    "Türkiye'nin Rüzgar Mirasından Geleceğin Teknolojisine",
    image_name="enercon_legacy_genesis_1778089178519.png",
    is_title_slide=True
)

# 2. Tarihsel Kilometre Taşları
add_magnificent_slide(
    "Türkiye'nin Rüzgar Serüvenindeki İlkler", 
    "Demirer Vizyonu ve Enercon Gücü",
    [
        "1996: Türkiye'nin ilk rüzgar ölçüm altyapısı (108 direk).",
        "1998: Germiyan RES - Türkiye'nin İLK rüzgar santrali.",
        "2000: Bozcaada RES - Enercon E-40 teknolojisi ile devrim.",
        "25+ Yıllık Kesintisiz İşletme Tecrübesi."
    ]
)

# 3. Neden Enercon? Teknolojik Üstünlük
add_magnificent_slide(
    "Enercon: Gearless (Dişlisiz) Teknolojinin Gücü", 
    "Minimum Bakım, Maksimum Verimlilik",
    [
        "Direct Drive Teknolojisi: Şanzıman kaynaklı arızaların eliminasyonu.",
        "Yüksek Emre Amadelik: Sektör standartlarının üzerinde performans.",
        "Şebeke Dostu: Dinamik reaktif güç kontrolü.",
        "Aerodinamik Mükemmellik: Sessiz ve verimli kanat tasarımları."
    ],
    image_name="enercon_direct_drive_tech_1778089196015.png"
)

# 4. Mevcut Filo Analizi
add_magnificent_slide(
    "Demirer Enerji Mevcut Filo Portföyü", 
    "284 Türbinlik Dev Enercon Ağı",
    [
        "Toplam Kurulu Güç: ~368 MW",
        "Toplam Enercon WEC: 284 Adet",
        "Ana Projeler: Çamseki, İntepe, Sayalar, Mare, Datça.",
        "Türkiye'nin en deneyimli rüzgar işletme kadrosu."
    ],
    image_name="fleet_analysis_infographic_1778088430993.png"
)

# 5. Kış Operasyonları ve Güvenlik
add_magnificent_slide(
    "Kış Operasyonları ve YAT Yönetimi", 
    "Ekstrem Koşullarda Operasyonel Süreklilik",
    [
        "Yük Atma (YAT) Talimatlarının donma riskine etkisi.",
        "Enercon 'Buz Algılama' ve 'Buz Çözme' sistemleri.",
        "Sensör verileriyle akıllı duruş ve güvenli restart.",
        "Operasyonel ve ticari risklerin minimize edilmesi."
    ],
    image_name="enercon_winter_grid_safety_1778089233070.png"
)

# 6. Repowering: Geleceğin Enerji Dönüşümü
add_magnificent_slide(
    "Repowering: Kapasite Artışı ve Modernizasyon", 
    "Eski Nesilden Yeni Nesil Enercon Teknolojisine",
    [
        "Mevcut sahalarda verimliliğin 2-3 kat artırılması.",
        "Yeni nesil Enercon E-138 / E-160 serisine geçiş vizyonu.",
        "Çevresel ayak izinin küçültülmesi, enerji üretiminin maksimize edilmesi."
    ],
    image_name="enercon_evolution_repowering_1778089214937.png"
)

# 7. Sonuç ve Vizyon
add_magnificent_slide(
    "Gelecek Rüzgarla Şekilleniyor", 
    "Demirer Holding & Enercon Ortaklığıyla",
    [
        "Sürdürülebilir büyüme stratejileri.",
        "Dijital dönüşüm ve siber güvenlik odaklı O&M.",
        "Türkiye'nin enerji bağımsızlığına katkı.",
        "Teşekkür ederiz. Soru & Cevap."
    ]
)

prs.save(output_file)
print(f"Magnificent PowerPoint file created at: {output_file}")

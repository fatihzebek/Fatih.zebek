from pptx import Presentation
import sys
import glob

files = glob.glob('*METRISO*.ppt*')
prs = Presentation(files[0])
for i, slide in enumerate(prs.slides):
    if i < 9: continue
    print(f'\\n--- Slide {i+1} ---')
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            # Convert to ascii, ignore weird chars
            print(shape.text.encode('ascii', 'ignore').decode('ascii'))
